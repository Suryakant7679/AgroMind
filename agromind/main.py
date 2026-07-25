import io
import base64
import time
import hashlib
import hmac
import os
import secrets
from pathlib import Path
from urllib.parse import quote

import markdown
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pptx import Presentation
from pptx.util import Inches, Pt
from starlette.middleware.sessions import SessionMiddleware

from agromind.ai import AIProviderError, generate_ai_response
from agromind.billing import all_plans, can_use_plan, estimate_cost, estimate_prompt_tokens, get_plan
from agromind.chatbot import WorkspaceChatRequest, WorkspaceChatbot
from agromind.data import DOMAINS, all_tools, get_domain, get_tool
from agromind.models import (
    DEFAULT_LANGUAGE,
    default_groq_model,
    get_language,
    groq_model_groups,
    groq_tts_model_for_language,
    language_options,
    resolve_response_language,
)
from agromind.supabase_store import (
    create_confirmed_user_with_password,
    fetch_profile,
    fetch_profiles,
    fetch_recent_outputs,
    fetch_output_by_id,
    _parse_time,
    insert_profile_if_missing,
    save_payment,
    save_output,
    save_subscription,
    sign_in_with_password,
    supabase_auth_configured,
    supabase_client,
    supabase_database_configured,
    update_profile_plan,
    update_profile,
    update_user_password,
    usage_summary,
    verify_email_token_hash,
    verify_signup_otp,
)

load_dotenv(".env.local")
load_dotenv()

# Clean environment variables of spaces/quotes
for env_key in [
    "NEXT_PUBLIC_SUPABASE_URL",
    "SUPABASE_URL",
    "NEXT_PUBLIC_SUPABASE_ANON_KEY",
    "SUPABASE_ANON_KEY",
    "SUPABASE_SERVICE_ROLE_KEY"
]:
    env_val = os.getenv(env_key)
    if env_val:
        os.environ[env_key] = env_val.strip().strip("'\"")

SERVER_START_TIME = time.time()
DEFAULT_DEV_SECRET = "agromind-dev-secret"
BASE_DIR = Path(__file__).resolve().parent
session_secret = os.getenv("SESSION_SECRET", DEFAULT_DEV_SECRET)
if os.getenv("ENVIRONMENT", "").lower() == "production" and session_secret == DEFAULT_DEV_SECRET:
    raise RuntimeError("Set a strong SESSION_SECRET before running in production.")

app = FastAPI(title="AgroMind AI")
app.add_middleware(SessionMiddleware, secret_key=session_secret, https_only=os.getenv("ENVIRONMENT", "").lower() == "production")
app.mount("/static", StaticFiles(directory=BASE_DIR / "static"), name="static")
templates = Jinja2Templates(directory=BASE_DIR / "templates")
workspace_chatbot = WorkspaceChatbot()

# Configure and ensure secure upload directory exists (handle read-only serverless filesystems gracefully)
UPLOADS_DIR = BASE_DIR / "static" / "uploads"
IS_READ_ONLY_FS = False
try:
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
except Exception as e:
    print(f"Warning: Upload directory could not be created ({e}). Falling back to /tmp.")
    IS_READ_ONLY_FS = True




@app.get("/api/health")
def health_check():
    return {"ok": True, "service": "agromind"}


@app.post("/api/chatbot/chat")
async def workspace_chat(request: Request, payload: WorkspaceChatRequest):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    return await workspace_chatbot.reply(
        payload,
        user_id=user_or_response.get("id"),
        access_token=user_or_response.get("access_token"),
    )


def user_from_session(request: Request) -> dict | None:
    return request.session.get("user")


def csrf_token(request: Request) -> str:
    token = request.session.get("csrf_token")
    if not token:
        token = secrets.token_urlsafe(32)
        request.session["csrf_token"] = token
    return token


def verify_csrf(request: Request, token: str | None) -> None:
    expected = request.session.get("csrf_token")
    if not expected or not token or not secrets.compare_digest(expected, token):
        raise HTTPException(status_code=403, detail="Invalid CSRF token.")


def safe_next(path: str) -> str:
    if not path.startswith("/") or path.startswith("//"):
        return "/dashboard"
    return path


def login_redirect(request: Request) -> RedirectResponse:
    return RedirectResponse(f"/login?next={quote(request.url.path)}", status_code=303)


def require_user(request: Request) -> dict | RedirectResponse:
    user = user_from_session(request)
    if not user:
        return login_redirect(request)
    return user


def require_admin(request: Request) -> tuple[dict, dict] | RedirectResponse:
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response

    profile = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token"))
    if not profile or profile.get("role") != "admin":
        return RedirectResponse("/dashboard", status_code=303)
    return user_or_response, profile


def page(request: Request, template: str, **context):
    base = {
        "request": request,
        "domains": DOMAINS,
        "user": user_from_session(request),
        "csrf_token": csrf_token(request),
        "languages": language_options(),
        "default_language": DEFAULT_LANGUAGE,
        "groq_models": groq_model_groups(),
        "plans": all_plans(),
        "razorpay_key_id": os.getenv("RAZORPAY_KEY_ID", ""),
    }
    base.update(context)
    return templates.TemplateResponse(request, template, base)


@app.get("/", response_class=HTMLResponse)
def home(request: Request):
    return page(request, "home.html", tools=all_tools())


@app.get("/dashboard", response_class=HTMLResponse)
def dashboard(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    recent = fetch_recent_outputs(user_or_response.get("id"), access_token=user_or_response.get("access_token"))
    
    enriched_recent = []
    for item in recent:
        domain_id = item.get("domain")
        tool_id = item.get("tool")
        domain = get_domain(domain_id)
        tool = get_tool(domain_id, tool_id)
        
        raw_date = item.get("created_at")
        pretty_date = "Recent"
        if raw_date:
            parsed = _parse_time(raw_date)
            if parsed:
                pretty_date = parsed.strftime("%b %d, %Y - %I:%M %p")
        
        enriched_recent.append({
            "id": item.get("id"),
            "domain_id": domain_id,
            "tool_id": tool_id,
            "domain_name": domain["name"] if domain else domain_id.replace("-", " ").capitalize(),
            "tool_title": tool["title"] if tool else tool_id.replace("-", " ").capitalize(),
            "pretty_date": pretty_date,
            "tokens_used": item.get("tokens_used", 0)
        })
        
    profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
    plan = get_plan(profile_data.get("plan", "starter"))
    usage = usage_summary(user_or_response.get("id"), user_or_response.get("access_token"))
    return page(request, "dashboard.html", tools=all_tools(), recent=enriched_recent, usage=usage, current_plan=plan)


@app.get("/chatbot", response_class=HTMLResponse)
def chatbot_page(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    ai_tutor_url = os.getenv("AI_TUTOR_URL", "").strip().rstrip("/")
    return page(
        request,
        "chatbot.html",
        ai_tutor_url=ai_tutor_url,
        ai_tutor_enabled=bool(ai_tutor_url),
    )


@app.get("/dashboard/{domain_id}", response_class=HTMLResponse)
def domain_page(request: Request, domain_id: str):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    domain = get_domain(domain_id)
    if not domain:
        return RedirectResponse("/dashboard", status_code=303)
    return page(request, "domain.html", domain=domain)


@app.get("/dashboard/{domain_id}/{tool_id}", response_class=HTMLResponse)
def tool_page(request: Request, domain_id: str, tool_id: str):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    domain = get_domain(domain_id)
    tool = get_tool(domain_id, tool_id)
    if not domain or not tool:
        return RedirectResponse("/dashboard", status_code=303)
    return page(request, "tool.html", domain=domain, tool=tool, output_html=None, fields={}, output_id=None)


@app.post("/dashboard/{domain_id}/{tool_id}", response_class=HTMLResponse)
async def run_tool(
    request: Request,
    domain_id: str,
    tool_id: str,
    asset: UploadFile | None = File(default=None),
    asset_camera: UploadFile | None = File(default=None),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    domain = get_domain(domain_id)
    tool = get_tool(domain_id, tool_id)
    if not domain or not tool:
        return RedirectResponse("/dashboard", status_code=303)
    profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
    plan_id = profile_data.get("plan", "starter")
    current_usage = usage_summary(user_or_response.get("id"), user_or_response.get("access_token"))
    allowed, limit_error = can_use_plan(plan_id, current_usage)
    if not allowed:
        return page(
            request,
            "tool.html",
            domain=domain,
            tool=tool,
            output_html=None,
            fields={},
            error=limit_error,
            response_language=DEFAULT_LANGUAGE,
            usage=current_usage,
            current_plan=get_plan(plan_id),
        )

    form = await request.form()
    verify_csrf(request, str(form.get("csrf_token", "")))
    fields = {field["name"]: str(form.get(field["name"], "")) for field in tool["fields"]}
    language_code = str(form.get("__language", DEFAULT_LANGUAGE))
    uploaded_asset = asset_camera if asset_camera and asset_camera.filename else asset
    
    # Secure validation & saving of uploaded file
    if uploaded_asset and uploaded_asset.filename:
        try:
            # 1. Size Validation (Max 10MB)
            MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
            content = await uploaded_asset.read()
            file_size = len(content)
            await uploaded_asset.seek(0)  # Always seek back after read!
            
            if file_size > MAX_FILE_SIZE:
                raise ValueError("File is too large. Maximum size allowed is 10MB.")
                
            # 2. Type Validation
            content_type = uploaded_asset.content_type
            filename_lower = uploaded_asset.filename.lower()
            allowed_mimes = ["image/jpeg", "image/png", "image/gif", "image/webp", "application/pdf"]
            allowed_extensions = (".jpg", ".jpeg", ".png", ".gif", ".webp", ".pdf")
            
            is_valid_type = (content_type in allowed_mimes) or filename_lower.endswith(allowed_extensions)
            if not is_valid_type:
                raise ValueError("Invalid file type. Only standard images (JPG, PNG, GIF, WEBP) and PDF files are allowed.")
                
            # 3. Secure File Saving
            file_ext = os.path.splitext(uploaded_asset.filename)[1]
            if not file_ext and content_type:
                if "pdf" in content_type:
                    file_ext = ".pdf"
                elif "png" in content_type:
                    file_ext = ".png"
                elif "gif" in content_type:
                    file_ext = ".gif"
                elif "webp" in content_type:
                    file_ext = ".webp"
                else:
                    file_ext = ".jpg"
            
            unique_filename = f"{secrets.token_hex(8)}{file_ext}"
            
            # Save based on filesystem writeability (using /tmp fallback on read-only systems)
            if not IS_READ_ONLY_FS:
                try:
                    saved_path = UPLOADS_DIR / unique_filename
                    with open(saved_path, "wb") as f:
                        f.write(content)
                    print(f"[UPLOAD SUCCESS] {uploaded_asset.filename} -> {saved_path} ({file_size} bytes)")
                except Exception as e:
                    print(f"Failed to write to UPLOADS_DIR, falling back to /tmp: {e}")
                    saved_path = Path("/tmp") / unique_filename
                    with open(saved_path, "wb") as f:
                        f.write(content)
                    print(f"[UPLOAD SUCCESS TEMP] {uploaded_asset.filename} -> {saved_path} ({file_size} bytes)")
            else:
                saved_path = Path("/tmp") / unique_filename
                with open(saved_path, "wb") as f:
                    f.write(content)
                print(f"[UPLOAD SUCCESS TEMP] {uploaded_asset.filename} -> {saved_path} ({file_size} bytes)")
                
            # Reset file pointer again so downstream processes can read it fully
            await uploaded_asset.seek(0)
            
        except ValueError as exc:
            return page(
                request,
                "tool.html",
                domain=domain,
                tool=tool,
                output_html=None,
                fields=fields,
                error=str(exc),
                response_language=language_code,
            )
        except Exception as exc:
            return page(
                request,
                "tool.html",
                domain=domain,
                tool=tool,
                output_html=None,
                fields=fields,
                error=f"File saving failed: {str(exc)}",
                response_language=language_code,
            )

    try:
        input_tokens = estimate_prompt_tokens(fields)
        output, provider, response_language = await generate_ai_response(domain_id, tool_id, fields, uploaded_asset, language_code, plan_id)
        output_tokens = max(1, len(output) // 4)
        billing = estimate_cost(provider, input_tokens, output_tokens)
        saved_output_id = save_output(
            user_or_response.get("id"),
            domain_id,
            tool_id,
            fields,
            output,
            provider,
            input_tokens,
            billing["credits"],
            billing["cost_cents"],
            user_or_response.get("access_token"),
        )
        output_html = markdown.markdown(output, extensions=["tables", "fenced_code"])
        error = None
    except AIProviderError as exc:
        output_html = None
        response_language = language_code
        error = exc.user_message
        saved_output_id = None
    except Exception as exc:
        output_html = None
        response_language = language_code
        error = str(exc)
        saved_output_id = None
    return page(
        request,
        "tool.html",
        domain=domain,
        tool=tool,
        output_html=output_html,
        output_id=saved_output_id,
        fields=fields,
        error=error,
        response_language=response_language,
    )


@app.get("/analytics", response_class=HTMLResponse)
def analytics(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
    plan = get_plan(profile_data.get("plan", "starter"))
    usage = usage_summary(user_or_response.get("id"), user_or_response.get("access_token"))
    rows = [
        {
            "domain": domain,
            "requests": values["requests"],
            "tokens": values["tokens"],
            "credits": values["credits"],
        }
        for domain, values in sorted(usage["domain_rows"].items())
    ]
    provider_rows = [
        {"provider": provider, **values}
        for provider, values in sorted(usage["provider_rows"].items())
    ]
    return page(request, "analytics.html", rows=rows, provider_rows=provider_rows, usage=usage, current_plan=plan)


@app.get("/admin", response_class=HTMLResponse)
def admin(request: Request):
    admin_or_response = require_admin(request)
    if isinstance(admin_or_response, RedirectResponse):
        return admin_or_response
    users = fetch_profiles()
    return page(request, "admin.html", users=users)


@app.get("/profile", response_class=HTMLResponse)
def profile(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token"))
    usage = usage_summary(user_or_response.get("id"), user_or_response.get("access_token"))
    return page(request, "profile.html", profile=profile_data, usage=usage, current_plan=get_plan((profile_data or {}).get("plan")), success=None, error=None)


@app.post("/profile", response_class=HTMLResponse)
async def profile_update(
    request: Request,
    full_name: str = Form(...),
    organization: str = Form(""),
    role: str = Form(""),
    csrf_token: str = Form(...),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    verify_csrf(request, csrf_token)

    success_msg = None
    error_msg = None
    try:
        current_profile = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
        role_to_save = role
        if current_profile.get("role") == "admin":
            role_to_save = "admin"
        elif not role_to_save:
            role_to_save = "member"

        update_profile(
            user_or_response.get("id"),
            {
                "full_name": full_name,
                "organization": organization,
                "role": role_to_save,
            },
            user_or_response.get("access_token"),
        )
        success_msg = "Profile updated successfully!"
    except Exception as e:
        error_msg = f"Failed to update profile: {str(e)}"

    profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token"))
    usage = usage_summary(user_or_response.get("id"), user_or_response.get("access_token"))
    return page(
        request,
        "profile.html",
        profile=profile_data,
        usage=usage,
        current_plan=get_plan((profile_data or {}).get("plan")),
        success=success_msg,
        error=error_msg,
    )


@app.get("/settings", response_class=HTMLResponse)
def settings(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    configured = {
        "supabase_auth": supabase_auth_configured(),
        "supabase_database": supabase_database_configured(),
        "gemini": bool(os.getenv("GEMINI_API_KEY")),
        "groq": bool(os.getenv("GROQ_API_KEY")),
    }
    defaults = {
        "reasoning": default_groq_model("reasoning"),
        "text": default_groq_model("text"),
        "vision": default_groq_model("vision"),
        "multilingual": default_groq_model("multilingual"),
        "speech_to_text": default_groq_model("speech_to_text"),
        "text_to_speech": default_groq_model("text_to_speech"),
    }
    return page(request, "settings.html", configured=configured, groq_defaults=defaults)


@app.get("/pricing", response_class=HTMLResponse)
def pricing(request: Request):
    user = user_from_session(request)
    profile_data = fetch_profile(user.get("id"), user.get("access_token")) if user else None
    return page(request, "pricing.html", current_plan_id=(profile_data or {}).get("plan", "starter"))


@app.get("/login", response_class=HTMLResponse)
def login(request: Request, next: str = "/dashboard"):
    return page(request, "login.html", next=safe_next(next), message=None)


@app.post("/login", response_class=HTMLResponse)
async def login_submit(
    request: Request,
    email: str = Form(...),
    password: str = Form(""),
    csrf_token: str = Form(...),
    next: str = Form("/dashboard"),
):
    verify_csrf(request, csrf_token)
    next = safe_next(next)
    if supabase_auth_configured():
        if not password:
            return page(request, "login.html", next=next, message="Password is required.")
        try:
            request.session["user"] = sign_in_with_password(email, password)
            return RedirectResponse(next, status_code=303)
        except Exception as exc:
            error_msg = str(exc)
            if "email not confirmed" in error_msg.lower():
                error_msg = "Your email is not yet confirmed. Please check your inbox for the verification link."
            elif "invalid" in error_msg.lower():
                error_msg = "Invalid email or password. Please try again."
            else:
                error_msg = f"Login failed: {error_msg}"
            return page(request, "login.html", next=next, message=error_msg)

    if os.getenv("ALLOW_DEMO_LOGIN", "").lower() == "true":
        request.session["user"] = {"id": None, "email": email}
        return RedirectResponse(next, status_code=303)

    return page(
        request,
        "login.html",
        next=next,
        message="Authentication is not configured. Add valid Supabase URL and anon key before signing in.",
    )


@app.get("/signup", response_class=HTMLResponse)
def signup(request: Request, next: str = "/dashboard"):
    return page(request, "signup.html", next=safe_next(next), message=None)


def complete_confirmed_signup(request: Request, email: str, password: str, full_name: str, next: str):
    user = create_confirmed_user_with_password(email, password, full_name)
    insert_profile_if_missing(user["id"], user["email"], full_name)
    signed_in_user = sign_in_with_password(email, password)
    request.session["user"] = signed_in_user
    for key in ("signup_email", "signup_full_name", "signup_next", "signup_verification_method", "signup_otp_type", "signup_password"):
        request.session.pop(key, None)
    return RedirectResponse(next, status_code=303)


@app.post("/signup", response_class=HTMLResponse)
async def signup_submit(
    request: Request,
    full_name: str = Form(""),
    email: str = Form(...),
    password: str = Form(...),
    csrf_token: str = Form(...),
    next: str = Form("/dashboard"),
):
    verify_csrf(request, csrf_token)
    next = safe_next(next)
    if len(password) < 6:
        return page(request, "signup.html", next=next, message="Password must be at least 6 characters.")

    if not supabase_auth_configured():
        return page(request, "signup.html", next=next, message="Supabase auth is not configured.")

    try:
        return complete_confirmed_signup(request, email, password, full_name, next)
    except Exception as exc:
        return page(request, "signup.html", next=next, message=str(exc))


@app.get("/check-email", response_class=HTMLResponse)
def check_email_page(request: Request):
    email = request.session.get("signup_email", "")
    next_dest = request.session.get("signup_next", "/dashboard")
    return page(
        request,
        "verify_otp.html",
        email=email,
        next=next_dest,
        mode="link",
        message=None,
    )


@app.get("/verify-otp", response_class=HTMLResponse)
def verify_otp_page(request: Request):
    email = request.session.get("signup_email", "")
    next_dest = request.session.get("signup_next", "/dashboard")
    return page(request, "verify_otp.html", email=email, next=next_dest, mode="otp", message=None)


@app.post("/verify-otp", response_class=HTMLResponse)
async def verify_otp_submit(
    request: Request,
    email: str = Form(...),
    token: str = Form(...),
    csrf_token: str = Form(...),
    next: str = Form("/dashboard"),
):
    verify_csrf(request, csrf_token)
    next = safe_next(next)
    full_name = request.session.get("signup_full_name", "")
    otp_type = request.session.get("signup_otp_type", "signup")
    signup_password = request.session.get("signup_password", "")
    try:
        user = verify_signup_otp(email, token, otp_type)
        if user.get("id") and user.get("email"):
            # Set the password using their access token if they verified via OTP
            if signup_password and user.get("access_token"):
                update_user_password(user["access_token"], signup_password)

            insert_profile_if_missing(user["id"], user["email"], full_name)
            request.session["user"] = user
            
            # Clean up session
            for key in ("signup_email", "signup_full_name", "signup_next",
                        "signup_verification_method", "signup_otp_type", "signup_password"):
                request.session.pop(key, None)
                
            return RedirectResponse(next, status_code=303)
        else:
            raise RuntimeError("Verification succeeded but did not return a valid user.")
    except Exception as exc:
        return page(request, "verify_otp.html", email=email, next=next, mode="otp", message=str(exc))


@app.get("/auth/confirmed", response_class=HTMLResponse)
def auth_confirmed(
    request: Request,
    next: str = "/dashboard",
    error: str = "",
    error_description: str = "",
    token_hash: str = "",
    type: str = "",
):
    next = safe_next(next)
    if error or error_description:
        message = error_description or error or "The confirmation link could not be verified. Please try signing up again."
        return page(request, "login.html", next=next, message=message)

    # Exchange the token_hash from the confirmation link for a real session
    if token_hash and type:
        try:
            user = verify_email_token_hash(token_hash, type)
            if user.get("id") and user.get("email"):
                full_name = request.session.get("signup_full_name", "")
                insert_profile_if_missing(user["id"], user["email"], full_name)
                request.session["user"] = user
                for key in ("signup_email", "signup_full_name", "signup_next",
                            "signup_verification_method", "signup_otp_type"):
                    request.session.pop(key, None)
                return RedirectResponse(next, status_code=303)
        except Exception as exc:
            return page(
                request, "login.html", next=next,
                message=f"Could not verify email automatically: {exc}. Please log in with your password.",
            )

    # Fallback: no token params — ask user to log in manually
    request.session.pop("signup_email", None)
    request.session.pop("signup_full_name", None)
    request.session.pop("signup_next", None)
    request.session.pop("signup_verification_method", None)
    request.session.pop("signup_otp_type", None)
    return page(
        request,
        "login.html",
        next=next,
        message="Email confirmed successfully. Please log in with your email and password to continue.",
    )


@app.get("/logout")
def logout(request: Request):
    request.session.clear()
    return RedirectResponse("/", status_code=303)


@app.get("/api/usage")
def usage():
    return {"message": "Use authenticated /analytics for real user usage data."}


@app.get("/system-monitor", response_class=HTMLResponse)
def system_monitor(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        return user_or_response
    
    # Calculate uptime
    uptime_seconds = time.time() - SERVER_START_TIME
    days = int(uptime_seconds // (24 * 3600))
    hours = int((uptime_seconds % (24 * 3600)) // 3600)
    minutes = int((uptime_seconds % 3600) // 60)
    seconds = int(uptime_seconds % 60)
    
    uptime_parts = []
    if days > 0: uptime_parts.append(f"{days}d")
    if hours > 0: uptime_parts.append(f"{hours}h")
    if minutes > 0: uptime_parts.append(f"{minutes}m")
    uptime_parts.append(f"{seconds}s")
    uptime_str = " ".join(uptime_parts)

    # Validate environment variables
    env_status = {
        "GROQ_API_KEY": {
            "name": "GROQ_API_KEY",
            "type": "Required",
            "status": "configured" if os.getenv("GROQ_API_KEY") else "missing",
            "description": "Required for voice STT transcription and dynamic agriculture calculator reasoning."
        },
        "GEMINI_API_KEY": {
            "name": "GEMINI_API_KEY",
            "type": "Required",
            "status": "configured" if os.getenv("GEMINI_API_KEY") else "missing",
            "description": "Required for crop disease analysis image processing and fallback LLM support."
        },
        "DATA_GOV_IN_API_KEY": {
            "name": "DATA_GOV_IN_API_KEY",
            "type": "Optional",
            "status": "configured" if os.getenv("DATA_GOV_IN_API_KEY") else "missing (fallback active)",
            "description": "Optional government key. Falls back to static agronomical historical pricing database if missing."
        },
        "PLANT_ID_API_KEY": {
            "name": "PLANT_ID_API_KEY",
            "type": "Optional",
            "status": "configured" if os.getenv("PLANT_ID_API_KEY") else "missing (Gemini vision active)",
            "description": "Optional leaf analysis key. Falls back to highly-accurate Gemini vision diagnostic scanner if missing."
        },
        "SUPABASE_URL": {
            "name": "SUPABASE_URL",
            "type": "Required",
            "status": "configured" if os.getenv("SUPABASE_URL") else "missing",
            "description": "Required for farmer database sync, pricing plans, and dashboard profile audits."
        }
    }

    # Let's check if the user is admin
    is_admin = False
    try:
        profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
        if profile_data.get("role") == "admin":
            is_admin = True
    except Exception:
        pass

    # Fetch recent outputs
    recent = []
    try:
        if is_admin:
            recent = fetch_recent_outputs(limit=50, access_token=user_or_response.get("access_token"))
        else:
            recent = fetch_recent_outputs(user_id=user_or_response.get("id"), limit=50, access_token=user_or_response.get("access_token"))
    except Exception:
        pass

    # Fetch profiles to map user_id -> email/name
    profiles_map = {}
    try:
        if is_admin:
            all_users = fetch_profiles() or []
            for p in all_users:
                profiles_map[str(p.get("id"))] = p
        else:
            profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token")) or {}
            profiles_map[str(user_or_response.get("id"))] = profile_data
    except Exception:
        pass

    # Group outputs by user details tuple: (name, email)
    grouped_history = {}
    for item in recent:
        item_user_id = str(item.get("user_id"))
        user_info = profiles_map.get(item_user_id) or {}
        user_name = user_info.get("full_name") or user_info.get("email") or "AgroMind User"
        user_email = user_info.get("email") or "unknown@agromind.ai"
        
        domain_id = item.get("domain")
        tool_id = item.get("tool")
        domain = get_domain(domain_id)
        tool = get_tool(domain_id, tool_id)
        
        raw_date = item.get("created_at")
        pretty_date = "Recent"
        if raw_date:
            parsed = _parse_time(raw_date)
            if parsed:
                pretty_date = parsed.strftime("%b %d, %Y - %I:%M %p")
                
        enriched_item = {
            "id": item.get("id"),
            "domain_id": domain_id,
            "tool_id": tool_id,
            "domain_name": domain["name"] if domain else domain_id.replace("-", " ").capitalize(),
            "tool_title": tool["title"] if tool else tool_id.replace("-", " ").capitalize(),
            "pretty_date": pretty_date,
            "tokens_used": item.get("tokens_used", 0)
        }
        
        user_key = (user_name, user_email)
        grouped_history.setdefault(user_key, [])
        grouped_history[user_key].append(enriched_item)

    return page(request, "system_monitor.html", uptime=uptime_str, env_status=env_status, grouped_history=grouped_history)


@app.post("/api/health-check")
async def api_health_check(request: Request):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Authentication required.")
    
    form = await request.form()
    verify_csrf(request, str(form.get("csrf_token", "")))
    
    from agromind.health import run_all_health_checks
    try:
        results = await run_all_health_checks()
        return {"ok": True, "results": results}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Diagnostic check suite crashed: {str(exc)}")


@app.post("/api/billing/create-order")
async def create_billing_order(request: Request, plan: str = Form(...), csrf_token: str = Form(...)):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)

    plan_config = get_plan(plan)
    if plan_config["price_inr"] <= 0:
        raise HTTPException(status_code=400, detail="This plan does not require payment.")

    merchant_upi = os.getenv("MERCHANT_UPI_ID", "surya@okaxis")
    merchant_name = os.getenv("MERCHANT_NAME", "AgroMind AI")
    amount = int(plan_config["price_inr"] * 100)
    
    order_id = f"upipay-{secrets.token_hex(4)}"
    
    save_payment(user_or_response.get("id"), plan, amount, order_id, status="created")
    return {
        "key_id": "upi",
        "order_id": order_id,
        "amount": amount,
        "currency": "INR",
        "plan_name": plan_config["name"],
        "upi_id": merchant_upi,
        "merchant_name": merchant_name,
    }


@app.post("/api/billing/verify")
async def verify_billing_payment(
    request: Request,
    plan: str = Form(...),
    csrf_token: str = Form(...),
    razorpay_order_id: str = Form(...),
    razorpay_payment_id: str = Form(...),
    razorpay_signature: str = Form(""),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)

    if not razorpay_order_id.startswith("upipay-"):
        raise HTTPException(status_code=400, detail="Invalid UPI transaction order.")

    utr = (razorpay_payment_id or "").strip()
    if len(utr) != 12 or not utr.isdigit():
        raise HTTPException(
            status_code=422,
            detail="Invalid UPI Ref No (UTR)! Must be exactly a 12-digit numeric code."
        )

    plan_config = get_plan(plan)
    update_profile_plan(user_or_response.get("id"), plan, user_or_response.get("access_token"))
    save_subscription(user_or_response.get("id"), plan, user_or_response.get("access_token"))
    save_payment(
        user_or_response.get("id"),
        plan,
        int(plan_config["price_inr"] * 100),
        razorpay_order_id,
        utr,
        "paid",
    )
    return {"ok": True, "plan": plan}


@app.post("/api/billing/downgrade")
async def downgrade_billing_plan(request: Request, plan: str = Form(...), csrf_token: str = Form(...)):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)

    if plan.lower() != "starter":
        raise HTTPException(status_code=400, detail="Only downgrades to starter are supported via this route.")

    update_profile_plan(user_or_response.get("id"), "starter", user_or_response.get("access_token"))
    save_subscription(user_or_response.get("id"), "starter", user_or_response.get("access_token"))
    return {"ok": True, "plan": "starter"}


@app.post("/api/voice-assistant")
async def voice_assistant(
    request: Request,
    query: str = Form(...),
    csrf_token: str = Form(...),
    language: str = Form(DEFAULT_LANGUAGE),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)

    selected_language = resolve_response_language(query, language)

    # 1. Personalize with the user's name from database profile or session/email
    user_name = "User"
    email = user_or_response.get("email", "")
    try:
        profile_data = fetch_profile(user_or_response.get("id"), user_or_response.get("access_token"))
        if profile_data and profile_data.get("full_name"):
            user_name = profile_data.get("full_name").split(" ")[0]
    except Exception:
        pass
    if user_name == "User" and email:
        user_name = email.split("@")[0].capitalize()

    # 2. Retrieve conversational history from session (short-term memory)
    chat_history = request.session.get("voice_chat_history", [])
    # Limit to the last 10 messages (5 exchange pairs) to strictly prevent cookie bloat
    chat_history = chat_history[-10:]

    # 3. Create context-aware, personalized, and language-adaptive system prompt
    system_prompt = (
        f"You are the intelligent, conversational voice and chat assistant for AgroMind, named AgroMind AI. "
        f"The current user's name is {user_name}. "
        "AgroMind is a multi-domain AI platform for agriculture, healthcare, and education. "
        "Your available capabilities and dashboard tools are:\n"
        "- Agriculture AI: Crop recommendation/planning, plant health inspector (photos of plants), farm operations, market intelligence.\n"
        "- Healthcare AI: Symptom checker, wellness coaching, medical report analysis, medicine safety support, skin analyzer.\n"
        "- Education AI: Notes, worksheets, MCQs, tutoring, grading essays, and YouTube study tools (extracts transcripts from YouTube links to make summary notes).\n\n"
        "Directives:\n"
        "1. Keep your reply extremely brief, friendly, highly conversational, and direct (max 2-3 sentences) because it will be spoken out loud by text-to-speech.\n"
        "2. If the user asks you to perform a complex action (e.g., analyze a plant disease picture, check skin disease, summarize a youtube video, or grade an essay), "
        "politely guide them to navigate to that specific dedicated tool in their AgroMind Dashboard.\n"
        f"3. Speak directly to {user_name} when appropriate.\n"
        f"4. CRITICAL: Always respond in the EXACT same language as the user's prompt (even if the user queries in romanized/mixed script like Hinglish or Spanglish). "
        f"Align perfectly with their tone and language. Default to {selected_language['name']} if the query language is ambiguous."
    )

    provider = "unknown"
    try:
        if os.getenv("GROQ_API_KEY"):
            provider = "groq"
            from openai import OpenAI
            client = OpenAI(
                api_key=os.getenv("GROQ_API_KEY"),
                base_url="https://api.groq.com/openai/v1",
                timeout=45,
            )
            messages = [{"role": "system", "content": system_prompt}]
            for msg in chat_history:
                messages.append({"role": msg["role"], "content": msg["content"]})
            messages.append({"role": "user", "content": query})

            # Dynamic Model Selection for Voice Assistant
            from agromind.ai import REASONING_TERMS
            query_lower = query.lower()
            is_simple_voice = len(query.split()) < 12 and not any(term in query_lower for term in REASONING_TERMS)
            assistant_model = "llama-3.1-8b-instant" if is_simple_voice else default_groq_model("multilingual")

            completion = client.chat.completions.create(
                model=assistant_model,
                messages=messages,
            )
            response_text = completion.choices[0].message.content or ""
        elif os.getenv("GEMINI_API_KEY"):
            provider = "gemini"
            import google.generativeai as genai
            genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
            from agromind.ai import gemini_model_candidates
            model_name = gemini_model_candidates()[0]
            model = genai.GenerativeModel(model_name)
            
            # Format conversational history as natural dialogue for Gemini's text mode
            payload = f"{system_prompt}\n\n"
            for msg in chat_history:
                role = "User" if msg["role"] == "user" else "Assistant"
                payload += f"{role}: {msg['content']}\n"
            payload += f"User: {query}\nAssistant:"

            result = model.generate_content(payload, request_options={"timeout": 45})
            response_text = result.text or ""
        else:
            response_text = "I am ready to help, but no AI providers are currently configured. Please configure your API keys."
    except Exception as exc:
        response_text = f"Sorry, I encountered an error while processing your request: {str(exc)}"

    # 3.5 Record AI usage metrics
    if provider != "unknown" and not response_text.startswith("Sorry, I encountered an error"):
        try:
            from agromind.billing import estimate_text_tokens, estimate_cost
            history_text = " ".join(msg["content"] for msg in chat_history)
            full_prompt = f"{system_prompt} {history_text} {query}"
            input_tokens = estimate_text_tokens(full_prompt)
            output_tokens = max(1, len(response_text) // 4)
            billing = estimate_cost(provider, input_tokens, output_tokens)
            save_output(
                user_or_response.get("id"),
                "assistant",
                "chat",
                {"query": query},
                response_text,
                provider,
                input_tokens,
                billing["credits"],
                billing["cost_cents"],
                user_or_response.get("access_token"),
            )
        except Exception as db_exc:
            print(f"Error saving assistant usage to database: {db_exc}")

    # 4. Record interaction in short-term session chat history
    chat_history.append({"role": "user", "content": query})
    chat_history.append({"role": "assistant", "content": response_text})
    request.session["voice_chat_history"] = chat_history[-10:]

    return {"response": response_text, "language": selected_language["code"]}


@app.post("/api/speech-to-text")
async def speech_to_text(
    request: Request,
    audio: UploadFile = File(...),
    csrf_token: str = Form(...),
    language: str = Form(DEFAULT_LANGUAGE),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)
    if not os.getenv("GROQ_API_KEY"):
        raise HTTPException(status_code=503, detail="Groq API key is not configured.")

    selected_language = get_language(language)
    content = await audio.read()
    if not content:
        raise HTTPException(status_code=400, detail="Audio file is empty.")

    try:
        from openai import OpenAI

        client = OpenAI(api_key=os.getenv("GROQ_API_KEY"), base_url="https://api.groq.com/openai/v1", timeout=60)
        model = default_groq_model("speech_to_text")
        try:
            transcript = client.audio.transcriptions.create(
                file=(audio.filename or "audio.webm", content, audio.content_type or "audio/webm"),
                model=model,
                language=selected_language["groq_code"],
            )
            return {"text": getattr(transcript, "text", ""), "model": model}
        except Exception as primary_exc:
            backup_model = "whisper-large-v3" if model == "whisper-large-v3-turbo" else "whisper-large-v3-turbo"
            print(f"Primary Whisper model {model} failed: {primary_exc}. Trying backup {backup_model}...")
            transcript = client.audio.transcriptions.create(
                file=(audio.filename or "audio.webm", content, audio.content_type or "audio/webm"),
                model=backup_model,
                language=selected_language["groq_code"],
            )
            return {"text": getattr(transcript, "text", ""), "model": backup_model}
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Groq speech-to-text failed: {exc}") from exc


@app.post("/api/text-to-speech")
async def text_to_speech(
    request: Request,
    text: str = Form(...),
    csrf_token: str = Form(...),
    language: str = Form(DEFAULT_LANGUAGE),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)
    if not os.getenv("GROQ_API_KEY"):
        raise HTTPException(status_code=503, detail="Groq API key is not configured.")

    selected_language = get_language(language)
    model = groq_tts_model_for_language(selected_language["code"])
    if not model:
        raise HTTPException(
            status_code=422,
            detail="Groq text-to-speech currently supports English and Arabic Saudi only. Browser TTS is used for the other supported app languages.",
        )
    try:
        from openai import OpenAI

        client = OpenAI(api_key=os.getenv("GROQ_API_KEY"), base_url="https://api.groq.com/openai/v1", timeout=60)
        speech = client.audio.speech.create(
            model=model,
            voice="tara",
            input=text[:4000],
            response_format="wav",
        )
        audio_bytes = speech.read() if hasattr(speech, "read") else bytes(speech.content)
        return StreamingResponse(io.BytesIO(audio_bytes), media_type="audio/wav", headers={"X-AgroMind-Model": model})
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Groq text-to-speech failed: {exc}") from exc


@app.get("/api/history/{output_id}")
async def get_history_detail(request: Request, output_id: str):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    
    output = fetch_output_by_id(output_id, user_or_response.get("access_token"))
    if not output:
        raise HTTPException(status_code=404, detail="Output not found.")
    
    # Check that this output belongs to the authenticated user
    if output.get("user_id") and str(output.get("user_id")) != str(user_or_response.get("id")):
        raise HTTPException(status_code=403, detail="Access denied.")
    
    # Resolve tool title and domain title
    tool_id = output.get("tool")
    domain_id = output.get("domain")
    domain = get_domain(domain_id)
    tool = get_tool(domain_id, tool_id)
    
    tool_title = tool["title"] if tool else tool_id.replace("-", " ").capitalize()
    domain_name = domain["name"] if domain else domain_id.replace("-", " ").capitalize()
    
    # Render markdown to HTML
    output_html = markdown.markdown(output.get("output_markdown", ""), extensions=["tables", "fenced_code"])
    
    return {
        "id": output.get("id"),
        "tool_id": tool_id,
        "domain_id": domain_id,
        "tool_title": tool_title,
        "domain_name": domain_name,
        "prompt": output.get("prompt", {}),
        "output_markdown": output.get("output_markdown", ""),
        "output_html": output_html,
        "tokens_used": output.get("tokens_used", 0),
        "created_at": output.get("created_at"),
    }


@app.post("/api/export/report")
async def export_report(request: Request, format: str = Form("pdf"), csrf_token: str = Form(...)):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)
    return {"ok": True, "format": format, "message": "Server-side export hook is ready in Python."}


@app.post("/api/ppt")
async def create_ppt(
    request: Request,
    title: str = Form("AgroMind AI Presentation"),
    slides: str = Form(""),
    csrf_token: str = Form(...),
):
    user_or_response = require_user(request)
    if isinstance(user_or_response, RedirectResponse):
        raise HTTPException(status_code=401, detail="Login required.")
    verify_csrf(request, csrf_token)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = title_slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(0.8))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True

    # Parse slide structures dynamically
    import json
    slide_data = []
    try:
        slide_data = json.loads(slides)
    except Exception:
        # Fallback to plain text splitting by newline
        slide_items = [item.strip() for item in slides.splitlines() if item.strip()] or ["Overview", "Key concepts", "Examples", "Summary"]
        for index, item in enumerate(slide_items, start=1):
            slide_data.append({
                "title": f"{index}. {item}",
                "content": ["Use this slide for key points, examples, and supporting notes from the selected topic."]
            })

    for item in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add heading
        heading = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.8))
        tf_heading = heading.text_frame
        tf_heading.word_wrap = True
        p_head = tf_heading.paragraphs[0]
        p_head.text = item.get("title", "Untitled Slide")
        p_head.font.size = Pt(28)
        p_head.font.bold = True
        
        # Add body content
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.8))
        tf_body = body.text_frame
        tf_body.word_wrap = True
        
        bullets = item.get("content", [])
        if isinstance(bullets, str):
            bullets = [bullets]
            
        for idx, bullet in enumerate(bullets):
            p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(12)
            p.font.size = Pt(18)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="agromind-presentation.pptx"'},
    )
